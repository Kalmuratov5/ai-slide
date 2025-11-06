import os
import logging
from telegram import Update
from telegram.ext import Application, CommandHandler, MessageHandler, filters, ContextTypes
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import google.generativeai as genai
from io import BytesIO

# 🔹 Logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# 🔹 Tokens va API kalitlari
TELEGRAM_BOT_TOKEN = os.getenv("TELEGRAM_BOT_TOKEN")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

# 🔹 Gemini AI konfiguratsiyasi
genai.configure(api_key=GEMINI_API_KEY)
model = genai.GenerativeModel("gemini-1.5")

# 🔹 /start komandasi
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "👋 Salom! Men Gemini AI botman.\n"
        "Mavzu yuboring, men sizga dizaynli slayd (.pptx) tayyorlab beraman.\n"
        "Masalan: 'O'zbekistonning iqlimi haqida slayd tayyorla'."
    )

# 🔹 Slayd tayyorlash funksiyasi
def create_pptx(title, slides_content):
    prs = Presentation()

    # Slayd dizayn va rang sozlamalari
    slide_layout = prs.slide_layouts[1]  # Title and Content
    colors = [
        RGBColor(79, 129, 189),
        RGBColor(192, 80, 77),
        RGBColor(155, 187, 89),
        RGBColor(128, 100, 162)
    ]

    for i, content in enumerate(slides_content):
        slide = prs.slides.add_slide(slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = colors[i % len(colors)]

        # Title
        title_placeholder = slide.shapes.title
        title_placeholder.text = content['title']
        title_placeholder.text_frame.paragraphs[0].font.size = Pt(44)
        title_placeholder.text_frame.paragraphs[0].font.bold = True
        title_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Content
        body = slide.placeholders[1]
        body.text = content['text']
        for paragraph in body.text_frame.paragraphs:
            paragraph.font.size = Pt(28)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
    
    # Faylni BytesIO ga saqlash
    pptx_file = BytesIO()
    prs.save(pptx_file)
    pptx_file.seek(0)
    return pptx_file

# 🔹 AI javobini olish va PPTX tayyorlash
async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    topic = update.message.text
    logger.info(f"Foydalanuvchi mavzuni yubordi: {topic}")

    prompt = (
        f"Sizga mavzu berildi: {topic}\n"
        "Iltimos, bu mavzuda 3-5 bo‘limdan iborat qisqa va to‘liq slayd tayyorlang. "
        "Har bir bo‘limni title va content ko‘rinishida ajrating. JSON formatida chiqaring: "
        "{'title': 'sarlavha', 'text': 'mazmun'}"
    )

    try:
        response = model.generate_text(prompt)
        ai_text = response.text if hasattr(response, "text") else ""
        
        # JSON formatini tozalash
        import json
        slides_content = json.loads(ai_text.replace("'", '"'))  # AI ko‘pincha ' ishlatadi

        pptx_file = create_pptx(topic, slides_content)
        await update.message.reply_document(document=pptx_file, filename=f"{topic}.pptx")
    except Exception as e:
        logger.error(f"Xato: {e}")
        await update.message.reply_text("❌ Xatolik yuz berdi. Iltimos, keyinroq urinib ko‘ring.")

# 🔹 Botni ishga tushirish
def main():
    application = Application.builder().token(TELEGRAM_BOT_TOKEN).build()
    
    application.add_handler(CommandHandler("start", start))
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_message))

    logger.info("✅ Bot ishga tushdi... (Gemini AI bilan, PPTX slaydlar tayyorlaydi)")
    application.run_polling()

if __name__ == "__main__":
    main()
